from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Iterable

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_FILL_LUMA_THRESHOLD = 48
BRIGHT_BORDER_LUMA_THRESHOLD = 200


def _remove_shape(shape) -> None:
    element = shape.element
    element.getparent().remove(element)


def _iter_shapes_recursive(shapes) -> Iterable:
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes_recursive(shape.shapes)


def _shape_has_image(shape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.PICTURE


def _set_background_white(container) -> None:
    fill = container.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def _is_dark_rgb(rgb: RGBColor) -> bool:
    # Perceived luminance weighting in sRGB space.
    luma = 0.2126 * rgb[0] + 0.7152 * rgb[1] + 0.0722 * rgb[2]
    return luma <= DARK_FILL_LUMA_THRESHOLD


def _is_bright_rgb(rgb: RGBColor) -> bool:
    luma = 0.2126 * rgb[0] + 0.7152 * rgb[1] + 0.0722 * rgb[2]
    return luma >= BRIGHT_BORDER_LUMA_THRESHOLD


def _set_fill_white_if_dark(fill, force_if_unknown: bool = False) -> None:
    if fill is None:
        return

    # Some fill proxies (e.g. _NoneFill) raise when foreground color is accessed.
    if getattr(fill, "type", None) is None:
        if force_if_unknown:
            fill.solid()
            fill.fore_color.rgb = WHITE
        return

    try:
        fore_color = fill.fore_color
        rgb = fore_color.rgb
    except (AttributeError, TypeError, ValueError):
        if force_if_unknown:
            fill.solid()
            fill.fore_color.rgb = WHITE
        return

    if rgb is None:
        if force_if_unknown:
            fill.solid()
            fill.fore_color.rgb = WHITE
        return

    if _is_dark_rgb(rgb):
        fill.solid()
        fill.fore_color.rgb = WHITE


def _line_is_explicitly_missing(line) -> bool:
    ln = getattr(line, "_ln", None)
    if ln is None:
        return True

    try:
        return ln.find(".//a:noFill", namespaces=ln.nsmap) is not None
    except (AttributeError, TypeError, ValueError):
        return False


def _set_line_black_if_bright(line, force_if_unknown: bool = False) -> None:
    if line is None:
        return

    if _line_is_explicitly_missing(line):
        return

    fill = getattr(line, "fill", None)
    if fill is None:
        return

    if getattr(fill, "type", None) is None:
        if force_if_unknown:
            fill.solid()
            fill.fore_color.rgb = BLACK
        return

    try:
        rgb = fill.fore_color.rgb
    except (AttributeError, TypeError, ValueError):
        if force_if_unknown:
            fill.solid()
            fill.fore_color.rgb = BLACK
        return

    if rgb is None:
        if force_if_unknown:
            fill.solid()
            fill.fore_color.rgb = BLACK
        return

    if _is_bright_rgb(rgb):
        fill.solid()
        fill.fore_color.rgb = BLACK


def _set_text_frame_black(text_frame) -> None:
    for paragraph in text_frame.paragraphs:
        paragraph.font.color.rgb = BLACK
        for run in paragraph.runs:
            run.font.color.rgb = BLACK


def _set_shape_text_black(shape) -> None:
    has_text_frame = getattr(shape, "has_text_frame", False)
    has_table = getattr(shape, "has_table", False)
    force_unknown_colors = not _shape_has_image(shape)
    shape_fill = getattr(shape, "fill", None)
    has_unknown_fill = getattr(shape_fill, "type", None) is None

    _set_line_black_if_bright(
        getattr(shape, "line", None),
        force_if_unknown=force_unknown_colors and has_unknown_fill,
    )

    _set_fill_white_if_dark(
        getattr(shape, "fill", None),
        force_if_unknown=force_unknown_colors or has_text_frame or has_table,
    )

    if has_text_frame:
        _set_text_frame_black(shape.text_frame)

    if has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                _set_fill_white_if_dark(cell.fill, force_if_unknown=True)
                _set_text_frame_black(cell.text_frame)

    try:
        child_shapes = shape.shapes
    except AttributeError:
        child_shapes = None

    if child_shapes is not None:
        for child_shape in child_shapes:
            _set_shape_text_black(child_shape)


def _strip_template_shapes(template_shapes) -> None:
    for shape in list(template_shapes):
        is_placeholder = getattr(shape, "is_placeholder", False)
        if _shape_has_image(shape) or not is_placeholder:
            _remove_shape(shape)


def _grayscale_blob(image_blob: bytes) -> bytes:
    with Image.open(BytesIO(image_blob)) as image:
        grayscale = image.convert("L").convert("RGB")
        out = BytesIO()
        grayscale.save(out, format="PNG")
        return out.getvalue()


def _replace_picture_with_grayscale(slide, picture_shape) -> None:
    gray_blob = _grayscale_blob(picture_shape.image.blob)
    stream = BytesIO(gray_blob)

    new_picture = slide.shapes.add_picture(
        stream,
        picture_shape.left,
        picture_shape.top,
        picture_shape.width,
        picture_shape.height,
    )

    new_picture.crop_left = picture_shape.crop_left
    new_picture.crop_right = picture_shape.crop_right
    new_picture.crop_top = picture_shape.crop_top
    new_picture.crop_bottom = picture_shape.crop_bottom
    new_picture.rotation = picture_shape.rotation

    _remove_shape(picture_shape)


def process_presentation(input_path: str | Path, grayscale_images: bool = False) -> Path:
    """Process a PPTX file and save a new '-wolayout' variant next to it."""
    input_path = Path(input_path)
    if not input_path.exists():
        raise FileNotFoundError(f"Input file not found: {input_path}")
    if input_path.suffix.lower() != ".pptx":
        raise ValueError("Input file must be a .pptx file")

    presentation = Presentation(str(input_path))

    for master in presentation.slide_masters:
        _set_background_white(master)
        _strip_template_shapes(master.shapes)
        for shape in master.shapes:
            _set_shape_text_black(shape)

    for layout in presentation.slide_layouts:
        _set_background_white(layout)
        _strip_template_shapes(layout.shapes)
        for shape in layout.shapes:
            _set_shape_text_black(shape)

    for slide in presentation.slides:
        _set_background_white(slide)

        picture_shapes = []
        for shape in slide.shapes:
            _set_shape_text_black(shape)

        if grayscale_images:
            for shape in _iter_shapes_recursive(slide.shapes):
                if _shape_has_image(shape):
                    picture_shapes.append(shape)

        for picture_shape in picture_shapes:
            _replace_picture_with_grayscale(slide, picture_shape)

    output_path = input_path.with_name(f"{input_path.stem}-wolayout{input_path.suffix}")
    presentation.save(str(output_path))
    return output_path
